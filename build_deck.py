"""
ChartRunner pitch deck — narrated 12-slide submission deck.
Output: PITCH-DECK.pptx in workspace folder.

Palette: Solana-native dark
  Primary bg:    #0D0D14
  Card bg:       #14141C
  Text white:    #FFFFFF
  Text dim:      #9AA0A6
  Solana purple: #9945FF
  Solana green:  #14F195
  Accent red:    #FF4D6D
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from copy import deepcopy

# ---- palette ----
BG          = RGBColor(0x0D, 0x0D, 0x14)
CARD        = RGBColor(0x14, 0x14, 0x1C)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
DIM         = RGBColor(0x9A, 0xA0, 0xA6)
PURPLE      = RGBColor(0x99, 0x45, 0xFF)
GREEN       = RGBColor(0x14, 0xF1, 0x95)
RED         = RGBColor(0xFF, 0x4D, 0x6D)
GOLD        = RGBColor(0xF0, 0xB9, 0x0B)
DIVIDER     = RGBColor(0x2A, 0x2A, 0x36)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
SW = prs.slide_width
SH = prs.slide_height
BLANK = prs.slide_layouts[6]

def set_bg(slide, rgb=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def add_rect(slide, left, top, width, height, fill=CARD, line=None, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(0.75)
    if radius:
        # adjust corner radius
        try:
            shp.adjustments[0] = 0.08
        except Exception:
            pass
    return shp

def add_text(slide, left, top, width, height, text, *, size=18, color=WHITE, bold=False,
             align=PP_ALIGN.LEFT, font='Calibri', italic=False):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    if isinstance(text, list):
        for i, run_def in enumerate(text):
            run = p.add_run() if i > 0 else p.runs[0] if p.runs else p.add_run()
            run.text = run_def['text']
            run.font.name = run_def.get('font', font)
            run.font.size = Pt(run_def.get('size', size))
            run.font.color.rgb = run_def.get('color', color)
            run.font.bold = run_def.get('bold', bold)
            run.font.italic = run_def.get('italic', italic)
    else:
        run = p.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.color.rgb = color
        run.font.bold = bold
        run.font.italic = italic
    return tb

def add_paragraphs(slide, left, top, width, height, paragraphs, *, default_size=14,
                   default_color=WHITE, default_font='Calibri', line_space=1.25):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = True
    for i, par in enumerate(paragraphs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = par.get('align', PP_ALIGN.LEFT)
        p.line_spacing = par.get('line_space', line_space)
        if 'space_after' in par:
            p.space_after = Pt(par['space_after'])
        runs = par.get('runs') or [{'text': par.get('text', '')}]
        for j, rd in enumerate(runs):
            run = p.add_run() if j > 0 else (p.runs[0] if p.runs else p.add_run())
            run.text = rd.get('text', '')
            run.font.name = rd.get('font', default_font)
            run.font.size = Pt(rd.get('size', default_size))
            run.font.color.rgb = rd.get('color', default_color)
            run.font.bold = rd.get('bold', False)
            run.font.italic = rd.get('italic', False)
    return tb

def set_notes(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text

def add_dot(slide, cx, cy, d, color=GREEN):
    shp = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx - d/2, cy - d/2, d, d)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    return shp

def add_pill(slide, left, top, text, color=PURPLE, w=Inches(1.4), h=Inches(0.34), text_color=WHITE):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shp.fill.solid(); shp.fill.fore_color.rgb = color
    shp.line.fill.background()
    try: shp.adjustments[0] = 0.5
    except Exception: pass
    tf = shp.text_frame
    tf.margin_left = Emu(0); tf.margin_right = Emu(0)
    tf.margin_top = Emu(0); tf.margin_bottom = Emu(0)
    tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    run = p.add_run(); run.text = text
    run.font.name = 'Calibri'; run.font.size = Pt(11); run.font.bold = True
    run.font.color.rgb = text_color
    return shp

def page_number(slide, n, total):
    add_text(slide, Inches(12.5), Inches(7.05), Inches(0.7), Inches(0.3),
             f"{n} / {total}", size=10, color=DIM, align=PP_ALIGN.RIGHT)

def footer_brand(slide):
    add_text(slide, Inches(0.5), Inches(7.05), Inches(4), Inches(0.3),
             [{'text': 'ChartRunner', 'bold': True, 'color': WHITE, 'size': 11},
              {'text': '  ·  Phase 0 MVP', 'color': DIM, 'size': 11}],
             font='Calibri')

TOTAL = 12

# ============================================================
# SLIDE 1 — Title
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

# accent corner
add_rect(s, Inches(0), Inches(0), Inches(0.18), SH, fill=PURPLE)

# Big title
add_text(s, Inches(0.9), Inches(2.1), Inches(11), Inches(1.4),
         "ChartRunner", size=88, color=WHITE, bold=True, font='Arial Black')

# tagline
add_text(s, Inches(0.9), Inches(3.4), Inches(11), Inches(0.7),
         "Trade the chart. Survive the upside-down.",
         size=24, color=GREEN, italic=True, font='Calibri')

# subtag
add_text(s, Inches(0.9), Inches(4.15), Inches(11), Inches(0.5),
         "Gamified trading SDK · Solana devnet · single-file MVP",
         size=16, color=DIM, font='Calibri')

# bottom strip
add_rect(s, Inches(0.9), Inches(6.3), Inches(11.5), Inches(0.05), fill=DIVIDER)
add_text(s, Inches(0.9), Inches(6.4), Inches(8), Inches(0.4),
         [{'text': 'Phase 0 — MVP shipped', 'color': GREEN, 'bold': True, 'size': 12},
          {'text': '  ·  v0.9 submission package', 'color': DIM, 'size': 12}],
         font='Calibri')
add_text(s, Inches(8.5), Inches(6.4), Inches(4), Inches(0.4),
         "ssjjul3.github.io/chartrunner",
         size=12, color=DIM, align=PP_ALIGN.RIGHT, font='Consolas')

set_notes(s,
"Speak: 'ChartRunner. Fortnite meets Space Invaders meets a trading chart. Every trade is a game move.'\n"
"Pause. Let it breathe. Don't read the slide aloud — talk past it.")

# ============================================================
# SLIDE 2 — Problem
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
         "The problem", size=14, color=DIM, bold=True, font='Calibri')
add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.5),
         "Trading apps are a hospital monitor.",
         size=44, color=WHITE, bold=True, font='Arial Black')

# Big stat block, left
add_rect(s, Inches(0.7), Inches(2.7), Inches(5.5), Inches(3.6), fill=CARD, radius=False)
add_text(s, Inches(0.95), Inches(2.95), Inches(5), Inches(2),
         "74%", size=130, color=RED, bold=True, font='Arial Black')
add_text(s, Inches(0.95), Inches(5.0), Inches(5), Inches(0.5),
         "of new retail traders quit",
         size=18, color=WHITE, bold=True, font='Calibri')
add_text(s, Inches(0.95), Inches(5.4), Inches(5), Inches(0.5),
         "within their first 90 days",
         size=18, color=WHITE, font='Calibri')
add_text(s, Inches(0.95), Inches(5.95), Inches(5), Inches(0.3),
         "Source: Brokerchooser, 2024", size=10, color=DIM, italic=True, font='Calibri')

# Right side — 3 reasons
add_text(s, Inches(6.8), Inches(2.7), Inches(6), Inches(0.5),
         "Why they quit", size=16, color=GREEN, bold=True, font='Calibri')

reasons = [
    ("Hospital-monitor UX", "12 panels, 100 numbers, an order ticket asking for real money before they understand what any of it means."),
    ("No skill loop", "YouTube and PDFs are passive. Paper trading is fake money on the same scary screen. Nothing transfers."),
    ("On-ramp is a cliff", "Skill required to start ≈ skill required to be profitable. Beginners can't get past minute one."),
]
y = 3.3
for title, body in reasons:
    add_dot(s, Inches(6.95), Inches(y + 0.16), Inches(0.18), GREEN)
    add_text(s, Inches(7.25), Inches(y), Inches(5.6), Inches(0.4),
             title, size=14, color=WHITE, bold=True, font='Calibri')
    add_text(s, Inches(7.25), Inches(y + 0.4), Inches(5.6), Inches(0.7),
             body, size=11, color=DIM, font='Calibri')
    y += 1.1

footer_brand(s); page_number(s, 2, TOTAL)

set_notes(s,
"Speak: 'Trading apps are a hospital monitor. 74% of new retail traders quit in 90 days. "
"Not because they're stupid — because the on-ramp is a cliff.'\n"
"Look at the audience. This is the universal pain. Don't soft-pedal the stat.")

# ============================================================
# SLIDE 3 — Solution
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
         "The solution", size=14, color=DIM, bold=True, font='Calibri')
add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.5),
         "Turn the cliff into a game.",
         size=44, color=WHITE, bold=True, font='Arial Black')

# 3-column layout
cols = [
    ("Game first", "Real Binance candles. Three avatar physics modes. Walk the chart, fall through the close, fight bears in the upside-down. Every input has visible feedback.", PURPLE),
    ("Real primitives", "Bracket. Ladder. OCO. Hedge. Radar. Rescue. Six trading primitives wired as in-game abilities — not toy versions.", GREEN),
    ("SDK is the bridge", "Every primitive routes through ChartRunnerSDK. The same SDK we plug into Solana devnet for live trades. What you practice is what graduates.", GOLD),
]
col_w = Inches(4.0); col_gap = Inches(0.27); col_x0 = Inches(0.7)
col_y = Inches(2.8); col_h = Inches(3.6)
for i, (h, b, c) in enumerate(cols):
    x = col_x0 + i*(col_w + col_gap)
    add_rect(s, x, col_y, col_w, col_h, fill=CARD)
    add_rect(s, x, col_y, col_w, Inches(0.15), fill=c)  # color top stripe
    add_text(s, x + Inches(0.4), col_y + Inches(0.45), col_w - Inches(0.8), Inches(0.6),
             h, size=22, color=c, bold=True, font='Arial Black')
    add_text(s, x + Inches(0.4), col_y + Inches(1.15), col_w - Inches(0.8), Inches(2.4),
             b, size=13, color=WHITE, font='Calibri')

# Bottom callout
add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.4),
         [{'text': 'The constitutional rule: ', 'color': DIM, 'italic': True, 'size': 12},
          {'text': 'abilities never touch the canvas. The SDK is the only thing that issues orders.',
           'color': WHITE, 'size': 12, 'bold': True}],
         font='Calibri')

footer_brand(s); page_number(s, 3, TOTAL)

set_notes(s,
"Speak: 'We turn the cliff into a game. Every primitive — bracket, ladder, OCO — is taught as a mechanic. "
"The hand learns where the stop goes. And here's the trick: every mechanic in the game routes through a real SDK. "
"What you practice is what graduates.'\n"
"Point at the third column. That's the unlock.")

# ============================================================
# SLIDE 4 — The MVP (live demo cue)
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
         "What's built", size=14, color=DIM, bold=True, font='Calibri')
add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.5),
         "Phase 0 MVP — playable now.",
         size=42, color=WHITE, bold=True, font='Arial Black')

# Left big card — single file
add_rect(s, Inches(0.7), Inches(2.6), Inches(5.7), Inches(4.2), fill=CARD)
add_text(s, Inches(0.95), Inches(2.85), Inches(5.2), Inches(0.5),
         "Single file. Zero install.", size=18, color=GREEN, bold=True, font='Calibri')
add_text(s, Inches(0.95), Inches(3.4), Inches(5.2), Inches(0.6),
         "ChartRunner_Prototype.html", size=15, color=WHITE, bold=True, font='Consolas')
add_text(s, Inches(0.95), Inches(3.9), Inches(5.2), Inches(0.4),
         "open it in your browser. that's the install.",
         size=12, color=DIM, italic=True, font='Consolas')
add_rect(s, Inches(0.95), Inches(4.45), Inches(5.2), Inches(0.04), fill=DIVIDER)

stats = [
    ("~14,700", "lines of vanilla JS"),
    ("0", "external dependencies"),
    ("0", "build steps"),
    ("230+", "atomic version commits"),
]
sx = 0.95
for i, (n, lbl) in enumerate(stats):
    cx = Inches(sx + (i % 2) * 2.5)
    cy = Inches(4.7 + (i // 2) * 0.95)
    add_text(s, cx, cy, Inches(2.3), Inches(0.5),
             n, size=24, color=WHITE, bold=True, font='Arial Black')
    add_text(s, cx, cy + Inches(0.5), Inches(2.3), Inches(0.3),
             lbl, size=10, color=DIM, font='Calibri')

# Right side — feature checklist
add_text(s, Inches(6.9), Inches(2.6), Inches(6.3), Inches(0.5),
         "What ships in Phase 0",
         size=16, color=GREEN, bold=True, font='Calibri')

features = [
    ("Three avatar physics", "Runner · Flight · Upside-down"),
    ("Six trading primitives", "Bracket · Ladder · OCO · Hedge · Radar · Rescue"),
    ("Two-anchor laser placement", "Click two candles to land any tool"),
    ("Workbench Pine builder", "Bots · Strategies · Indicators · Backtest · Apps"),
    ("Multi-tracker Terminal", "Darkflow · Hyper · Solana · CEX · Strategies"),
    ("Drag-to-desktop widgets", "Map / chat / pane → live preview"),
    ("Phone OS overlay", "Full-screen mobile mirror"),
    ("Save / load maps", "30-slot setup recall with thumbnails"),
]
y = 3.2
for h, b in features:
    add_text(s, Inches(6.9), Inches(y), Inches(0.25), Inches(0.3),
             "✓", size=14, color=GREEN, bold=True, font='Calibri')
    add_text(s, Inches(7.18), Inches(y), Inches(6.1), Inches(0.3),
             [{'text': h, 'color': WHITE, 'size': 12, 'bold': True},
              {'text': '  ' + b, 'color': DIM, 'size': 11}],
             font='Calibri')
    y += 0.4

footer_brand(s); page_number(s, 4, TOTAL)

set_notes(s,
"Speak: 'I'm going to show you the prototype. It's one HTML file. No build, no install. "
"Open it. Let's place a bracket.'\n"
"DEMO TIME — switch to the browser and run the 3-min flow from PITCH-DELIVERY.md. "
"Don't read this slide. Talk over the screen.")

# ============================================================
# SLIDE 5 — Architecture
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
         "How it works", size=14, color=DIM, bold=True, font='Calibri')
add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.5),
         "Abilities never touch the canvas.",
         size=40, color=WHITE, bold=True, font='Arial Black')

# Architecture stack diagram — 3 layers
def stack_layer(slide, top, label_top, body_top, label, body, color):
    add_rect(slide, Inches(2.8), top, Inches(7.7), Inches(1.4), fill=CARD)
    add_rect(slide, Inches(2.8), top, Inches(0.18), Inches(1.4), fill=color)
    add_text(slide, Inches(3.15), label_top, Inches(7.2), Inches(0.5),
             label, size=18, color=color, bold=True, font='Arial Black')
    add_text(slide, Inches(3.15), body_top, Inches(7.2), Inches(0.7),
             body, size=12, color=DIM, font='Calibri')

stack_layer(s, Inches(2.6), Inches(2.75), Inches(3.2), "ChartRunner UI",
            "canvas · player physics · monsters · particles · tools (visual only)", PURPLE)

stack_layer(s, Inches(4.15), Inches(4.3), Inches(4.75), "ChartRunnerSDK",
            "placeBracket · placeLadder · placeOCO · openHedge · radarScan · rescue · event bus", GREEN)

stack_layer(s, Inches(5.7), Inches(5.85), Inches(6.3), "Adapter (Phase 2)",
            "Solana devnet · paper-mode mirror · Phantom-Connect lifecycle", GOLD)

# Connecting arrows
for y in [Inches(4.0), Inches(5.55)]:
    add_rect(s, Inches(6.55), y, Inches(0.18), Inches(0.15), fill=WHITE)

# Side callouts
add_text(s, Inches(0.7), Inches(2.85), Inches(2.0), Inches(0.5),
         "Game", size=11, color=DIM, bold=True, font='Calibri', align=PP_ALIGN.RIGHT)
add_text(s, Inches(0.7), Inches(4.4), Inches(2.0), Inches(0.5),
         "Contract", size=11, color=DIM, bold=True, font='Calibri', align=PP_ALIGN.RIGHT)
add_text(s, Inches(0.7), Inches(5.95), Inches(2.0), Inches(0.5),
         "Venue", size=11, color=DIM, bold=True, font='Calibri', align=PP_ALIGN.RIGHT)

# Right — design rule
add_text(s, Inches(10.8), Inches(3.0), Inches(2.4), Inches(2.5),
         [{'text': 'Why this matters\n\n', 'color': WHITE, 'size': 13, 'bold': True},
          {'text': 'Phase 2 is a swap, not a rewrite. Same SDK call, new adapter. ',
           'color': DIM, 'size': 11},
          {'text': 'Players take their muscle memory with them.',
           'color': GREEN, 'size': 11, 'bold': True}],
         font='Calibri')

footer_brand(s); page_number(s, 5, TOTAL)

set_notes(s,
"Speak: 'The constitutional rule: abilities never touch the canvas, the SDK is the only thing that "
"issues orders. That sounds like a small thing. It's the whole bet. It's why Phase 2 is a swap.'")

# ============================================================
# SLIDE 6 — Why now
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
         "Why now", size=14, color=DIM, bold=True, font='Calibri')
add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.5),
         "Two structural shifts.",
         size=42, color=WHITE, bold=True, font='Arial Black')

shifts = [
    ("01", "Wallet UX flipped",
     "Phantom + Hyperliquid lowered the on-chain trading on-ramp from days to minutes. A 22-year-old can fund a wallet and place a perp in under 4 minutes. The infrastructure barrier is gone.",
     PURPLE),
    ("02", "Memecoin season trained 8M+ wallets",
     "Pump.fun, Bonk, and the post-2024 launch culture taught a generation to swap on-chain. They have wallets. They have liquidity. They don't have a way to learn anything beyond \"press buy.\"",
     GREEN),
]
y = 2.7
for num, head, body, color in shifts:
    add_rect(s, Inches(0.7), Inches(y), Inches(11.9), Inches(1.85), fill=CARD)
    add_text(s, Inches(0.95), Inches(y + 0.25), Inches(1.4), Inches(1.4),
             num, size=64, color=color, bold=True, font='Arial Black')
    add_text(s, Inches(2.4), Inches(y + 0.3), Inches(10), Inches(0.6),
             head, size=22, color=WHITE, bold=True, font='Arial Black')
    add_text(s, Inches(2.4), Inches(y + 0.95), Inches(10), Inches(0.85),
             body, size=12, color=DIM, font='Calibri')
    y += 2.05

# Bottom punch
add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         [{'text': 'The infra is here. The skill on-ramp is missing. ', 'color': DIM, 'size': 13},
          {'text': 'Nobody owns it.', 'color': GREEN, 'size': 13, 'bold': True}],
         font='Calibri')

footer_brand(s); page_number(s, 6, TOTAL)

set_notes(s,
"Speak: 'Hyperliquid + Phantom flipped the wallet UX. Memecoin season trained 8 million wallets. "
"The infra is here. The skill on-ramp is missing. Nobody owns it.'")

# ============================================================
# SLIDE 7 — Traction
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
         "Traction", size=14, color=DIM, bold=True, font='Calibri')
add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.5),
         "Shipping discipline you can audit.",
         size=38, color=WHITE, bold=True, font='Arial Black')

# Big number left
add_rect(s, Inches(0.7), Inches(2.7), Inches(5.0), Inches(3.6), fill=CARD)
add_text(s, Inches(0.95), Inches(3.05), Inches(4.5), Inches(1.6),
         "230+", size=100, color=GREEN, bold=True, font='Arial Black')
add_text(s, Inches(0.95), Inches(4.85), Inches(4.5), Inches(0.5),
         "atomic version commits", size=18, color=WHITE, bold=True, font='Calibri')
add_text(s, Inches(0.95), Inches(5.35), Inches(4.5), Inches(0.4),
         "each parse-validated", size=13, color=DIM, font='Calibri')
add_text(s, Inches(0.95), Inches(5.7), Inches(4.5), Inches(0.4),
         "each with reasoning in source", size=13, color=DIM, font='Calibri')

# Right column — milestones
add_text(s, Inches(6.3), Inches(2.7), Inches(6), Inches(0.5),
         "Milestones already met", size=16, color=GREEN, bold=True, font='Calibri')

ms = [
    ("Phase 0 MVP", "Single-file playable, 9 OS apps, 5 trackers, 6 SDK primitives", True),
    ("Phase 1 SDK skeleton", "ES modules, RiskManager parity tests, demo page", True),
    ("M5 Solana paper-mode", "Ed25519 wallet + session + approver + parity vectors", True),
    ("Constitutional rule", "Enforced across 230+ commits — SDK never bypassed", True),
    ("Public demo URL", "GitHub Pages folder ready (chartrunner-prototype/)", True),
    ("Pitch package", "Deck + 9 docs + video script + X launch kit", True),
]
y = 3.25
for h, b, done in ms:
    color = GREEN if done else DIM
    mark = "●" if done else "○"
    add_text(s, Inches(6.3), Inches(y), Inches(0.3), Inches(0.3),
             mark, size=14, color=color, bold=True, font='Calibri')
    add_text(s, Inches(6.65), Inches(y), Inches(6), Inches(0.3),
             h, size=13, color=WHITE, bold=True, font='Calibri')
    add_text(s, Inches(6.65), Inches(y + 0.3), Inches(6), Inches(0.3),
             b, size=10, color=DIM, font='Calibri')
    y += 0.55

footer_brand(s); page_number(s, 7, TOTAL)

set_notes(s,
"Speak: 'We're pre-launch. What we have is shipping discipline. 230+ atomic version commits. "
"A single-file prototype that's been parse-validated and playtested at every step.'\n"
"If asked about users: 'Zero today. We ship publicly with the next push and we have a 30-day target board.'")

# ============================================================
# SLIDE 8 — Competitive
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
         "Competitive edge", size=14, color=DIM, bold=True, font='Calibri')
add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.5),
         "Four wedges nobody else has.",
         size=38, color=WHITE, bold=True, font='Arial Black')

# 2×2 grid
items = [
    ("Gamified", "TradingView and Phantom aren't.\nLearn-by-doing, minute one.", PURPLE),
    ("Skill-building", "Bitget Quest and Coinbase Learn aren't.\nReps build muscle memory, not trivia.", GREEN),
    ("On-chain native", "TradingView Paper isn't.\nPhase 2 routes through Solana devnet.", GOLD),
    ("SDK-portable", "Every gamified competitor locks the player to one venue.\nWe drop on any chart.", RED),
]
gw, gh, gx0, gy0, gap = Inches(5.85), Inches(2.0), Inches(0.7), Inches(2.7), Inches(0.2)
for i, (h, b, c) in enumerate(items):
    col, row = i % 2, i // 2
    x = gx0 + col*(gw + gap); y = gy0 + row*(gh + gap)
    add_rect(s, x, y, gw, gh, fill=CARD)
    add_rect(s, x, y, gw, Inches(0.12), fill=c)
    add_text(s, x + Inches(0.4), y + Inches(0.35), gw - Inches(0.8), Inches(0.5),
             h, size=22, color=c, bold=True, font='Arial Black')
    add_text(s, x + Inches(0.4), y + Inches(1.0), gw - Inches(0.8), Inches(0.9),
             b, size=12, color=WHITE, font='Calibri')

# Bottom punch — placed above footer with breathing room
add_text(s, Inches(0.7), Inches(6.7), Inches(12), Inches(0.35),
         [{'text': 'We are a runtime artifact, not a SaaS. ', 'color': DIM, 'size': 12},
          {'text': 'Distribution asymmetry.', 'color': GREEN, 'size': 12, 'bold': True}],
         font='Calibri')

page_number(s, 8, TOTAL)

set_notes(s,
"Speak: 'Four wedges. Gamified — TV and Phantom aren't. Skill-building — Bitget Quest isn't. "
"On-chain native — TV Paper isn't. SDK-portable — every other gamified product locks to one venue. "
"We're a runtime artifact, not a SaaS. We can drop on any chart.'")

# ============================================================
# SLIDE 9 — Roadmap
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
         "Roadmap", size=14, color=DIM, bold=True, font='Calibri')
add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.5),
         "From single-file demo to live Solana trades.",
         size=32, color=WHITE, bold=True, font='Arial Black')

# Three horizontal phases
phases = [
    ("Phase 0", "SHIPPED", "Playable single-file MVP\n· Real Binance candles\n· 6 SDK primitives\n· 5 tracker views\n· Workbench Pine builder", GREEN),
    ("Phase 1", "ARCHITECTURE DONE", "ChartHost adapter layer\n· Drop UI on Dexscreener\n· Drop on TradingView\n· Drop on Birdeye\n· npm @chartrunner/sdk@0.1", GOLD),
    ("Phase 2", "PLANNED", "Live Solana devnet\n· Phantom-Connect\n· Real on-chain fills\n· P2P creator economy\n· Multi-venue adapters", PURPLE),
]
gw, gh, gx0, gy = Inches(4.0), Inches(4.0), Inches(0.7), Inches(2.5)
for i, (lbl, status, body, color) in enumerate(phases):
    x = gx0 + i*(gw + Inches(0.2))
    add_rect(s, x, gy, gw, gh, fill=CARD)
    add_rect(s, x, gy, gw, Inches(0.4), fill=color)
    add_text(s, x + Inches(0.3), gy + Inches(0.05), gw - Inches(0.6), Inches(0.3),
             lbl, size=16, color=BG, bold=True, font='Arial Black')
    add_text(s, x + Inches(0.3), gy + Inches(0.55), gw - Inches(0.6), Inches(0.4),
             status, size=12, color=color, bold=True, font='Calibri')
    add_text(s, x + Inches(0.3), gy + Inches(1.05), gw - Inches(0.6), Inches(2.8),
             body, size=12, color=WHITE, font='Calibri')

# Timeline at bottom
add_text(s, Inches(0.7), Inches(6.85), Inches(12), Inches(0.4),
         [{'text': 'Now', 'color': GREEN, 'bold': True, 'size': 12},
          {'text': '  →  90 days  →  ', 'color': DIM, 'size': 12},
          {'text': '180 days', 'color': PURPLE, 'bold': True, 'size': 12},
          {'text': '  →  Phase 2 mainnet target', 'color': DIM, 'size': 12}],
         font='Calibri')

footer_brand(s); page_number(s, 9, TOTAL)

set_notes(s,
"Speak: 'Phase 0 shipped. Phase 1 — drop the ChartRunner UI on top of Dexscreener, TradingView, "
"Birdeye. Anywhere with a candle. Phase 2 — wallet-connect and live Solana trades. "
"We're not asking you to imagine the architecture. It's already there.'")

# ============================================================
# SLIDE 10 — Business model
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

add_text(s, Inches(0.7), Inches(0.5), Inches(8), Inches(0.6),
         "Business model", size=14, color=DIM, bold=True, font='Calibri')
add_text(s, Inches(0.7), Inches(0.95), Inches(12), Inches(1.5),
         "Three revenue lanes, one runtime.",
         size=36, color=WHITE, bold=True, font='Arial Black')

streams = [
    ("P2P creator marketplace", "5% fee on every Workbench-built bot, strategy, and map sold for $SOL.", "Creator-economy lock-in. TradingView can't match without breaking SaaS."),
    ("Game shop", "Vehicles, skins, themes priced in $RUN earned in-run. Free cosmetic loop with paid premium tier.", "F2P-style monetization, crypto-native settlement."),
    ("White-label license", "Per-venue annual license for chart vendors and exchanges that want a gamified onboarding tab.", "Phantom · Drift · Hyperliquid · Birdeye are all candidate partners."),
]
y = 2.7
for h, body, why in streams:
    add_rect(s, Inches(0.7), Inches(y), Inches(11.9), Inches(1.3), fill=CARD)
    add_text(s, Inches(0.95), Inches(y + 0.18), Inches(11.5), Inches(0.5),
             h, size=18, color=GREEN, bold=True, font='Arial Black')
    add_text(s, Inches(0.95), Inches(y + 0.65), Inches(7.5), Inches(0.6),
             body, size=12, color=WHITE, font='Calibri')
    add_text(s, Inches(8.7), Inches(y + 0.65), Inches(3.7), Inches(0.6),
             why, size=11, color=DIM, italic=True, font='Calibri')
    y += 1.4

footer_brand(s); page_number(s, 10, TOTAL)

set_notes(s,
"Speak: 'Three revenue lanes. Marketplace fees on creator content. Game-shop $RUN sinks. "
"White-label license for chart vendors. The runtime artifact lets us serve all three from one codebase.'")

# ============================================================
# SLIDE 11 — The ask
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s, BG)

# big purple wash on left
add_rect(s, Inches(0), Inches(0), Inches(5.2), SH, fill=PURPLE)

add_text(s, Inches(0.5), Inches(1.2), Inches(4.6), Inches(0.5),
         "The ask", size=14, color=WHITE, bold=True, font='Calibri')
add_text(s, Inches(0.5), Inches(1.7), Inches(4.6), Inches(2.0),
         "Two open invites.", size=44, color=WHITE, bold=True, font='Arial Black')
add_text(s, Inches(0.5), Inches(4.6), Inches(4.6), Inches(2),
         "Phase 0 ships now.\nPhase 1 in 90 days with a partner.\nPhase 2 in 180 with seed.",
         size=15, color=WHITE, font='Calibri')

# right side — the two asks
x0 = Inches(5.7); w = Inches(7.1); y0 = Inches(1.3)
add_rect(s, x0, y0, w, Inches(2.5), fill=CARD)
add_text(s, x0 + Inches(0.4), y0 + Inches(0.3), Inches(0.5), Inches(0.6),
         "1", size=44, color=GREEN, bold=True, font='Arial Black')
add_text(s, x0 + Inches(1.1), y0 + Inches(0.4), w - Inches(1.5), Inches(0.5),
         "Devnet integration partner",
         size=22, color=GREEN, bold=True, font='Arial Black')
add_text(s, x0 + Inches(1.1), y0 + Inches(1.0), w - Inches(1.5), Inches(1.3),
         "Hyperliquid · Drift · Phoenix — first-class fit.\n"
         "We bring the on-ramp; you bring the venue.\n"
         "Co-signed integration spec, 90-day Phase 1 ship target.",
         size=12, color=WHITE, font='Calibri')

y1 = Inches(4.0)
add_rect(s, x0, y1, w, Inches(2.5), fill=CARD)
add_text(s, x0 + Inches(0.4), y1 + Inches(0.3), Inches(0.5), Inches(0.6),
         "2", size=44, color=GOLD, bold=True, font='Arial Black')
add_text(s, x0 + Inches(1.1), y1 + Inches(0.4), w - Inches(1.5), Inches(0.5),
         "Six-month seed",
         size=22, color=GOLD, bold=True, font='Arial Black')
add_text(s, x0 + Inches(1.1), y1 + Inches(1.0), w - Inches(1.5), Inches(1.3),
         "Staffs the Phase 1 SDK extraction + Phase 2 Solana adapter.\n"
         "1 Solana eng · 1 comms hire · 6 months runway.\n"
         "Solana foundation grants in flight as parallel track.",
         size=12, color=WHITE, font='Calibri')

# Custom footer for slide 11 — purple wash on the left makes the standard
# DIM grey unreadable, so use white/green tokens for contrast.
add_text(s, Inches(0.5), Inches(7.05), Inches(4), Inches(0.3),
         [{'text': 'ChartRunner', 'bold': True, 'color': WHITE, 'size': 11},
          {'text': '  ·  Phase 0 MVP', 'color': GREEN, 'size': 11}],
         font='Calibri')
page_number(s, 11, TOTAL)

set_notes(s,
"Speak: 'Two things. One — a devnet integration partner. Hyperliquid, Drift, or Phoenix would be the "
"strongest fits. Two — a six-month seed to staff Phase 1 and ship Phase 2 live.'\n"
"State this flat. Don't apologize. Don't hedge. Then sit with the silence.")

# ============================================================
# SLIDE 12 — Closing
# ============================================================
s = prs.slides.add_slide(BLANK)
set_bg(s)

# top tag
add_rect(s, Inches(0), Inches(0), SW, Inches(0.18), fill=GREEN)

add_text(s, Inches(0.7), Inches(2.2), Inches(12), Inches(1.5),
         "The next 10 million traders",
         size=44, color=DIM, bold=True, font='Arial Black')
add_text(s, Inches(0.7), Inches(3.0), Inches(12), Inches(1.5),
         "won't learn from PDFs.",
         size=44, color=WHITE, bold=True, font='Arial Black')

add_text(s, Inches(0.7), Inches(4.4), Inches(12), Inches(0.7),
         "They'll learn from a game.",
         size=32, color=GREEN, italic=True, bold=True, font='Arial Black')

add_text(s, Inches(0.7), Inches(5.3), Inches(12), Inches(0.5),
         "We're the on-ramp.",
         size=20, color=DIM, font='Calibri')

# Contacts
add_rect(s, Inches(0.7), Inches(6.1), Inches(11.9), Inches(0.04), fill=DIVIDER)

add_text(s, Inches(0.7), Inches(6.3), Inches(4.3), Inches(0.4),
         [{'text': 'github.com/ssjjul3/chartrunner', 'color': WHITE, 'size': 11, 'bold': True}],
         font='Consolas')
add_text(s, Inches(5.0), Inches(6.3), Inches(4), Inches(0.4),
         [{'text': 'ssjjul3.github.io/chartrunner', 'color': WHITE, 'size': 12, 'bold': True}],
         font='Consolas', align=PP_ALIGN.CENTER)
add_text(s, Inches(8.8), Inches(6.3), Inches(4), Inches(0.4),
         [{'text': 'jsg@julianroy.com', 'color': WHITE, 'size': 12, 'bold': True}],
         font='Calibri', align=PP_ALIGN.RIGHT)

add_text(s, Inches(0.7), Inches(6.65), Inches(12), Inches(0.3),
         "Repo  ·  Live demo  ·  Email",
         size=10, color=DIM, font='Calibri')

page_number(s, 12, TOTAL)

set_notes(s,
"Speak: 'The next 10 million traders won't learn from PDFs. They'll learn from a game. "
"We're the on-ramp.'\n"
"Pause. Then: 'Questions?' \n"
"Don't fill the silence. Wait for the question.")

# ---- write ----
out = "/sessions/busy-vigilant-allen/mnt/Trading Game/PITCH-DECK.pptx"
prs.save(out)
print(f"OK wrote {out}")
print(f"Slides: {len(prs.slides)}")
